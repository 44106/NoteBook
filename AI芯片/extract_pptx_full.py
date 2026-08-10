from __future__ import annotations

import json
import os
import re
import sys
import zipfile
from dataclasses import dataclass, asdict
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
}


@dataclass
class ShapeRecord:
    path: str
    shape_type: str
    name: str
    left: int | None
    top: int | None
    width: int | None
    height: int | None
    text: str
    table: list[list[str]] | None = None
    chart: dict[str, Any] | None = None
    alt_text: str | None = None
    rel_ids: list[str] | None = None


def clean_text(value: str) -> str:
    value = value.replace("\r", "\n")
    value = re.sub(r"[ \t]+\n", "\n", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def shape_alt_text(shape: Any) -> str | None:
    try:
        c_nv_pr = shape.element.xpath(".//p:cNvPr")
        vals: list[str] = []
        for node in c_nv_pr:
            for attr in ("name", "descr", "title"):
                v = node.get(attr)
                if v and v not in vals:
                    vals.append(v)
        text = " | ".join(vals).strip()
        return text or None
    except Exception:
        return None


def shape_rel_ids(shape: Any) -> list[str]:
    try:
        ids = []
        for node in shape.element.xpath(".//*[@r:embed or @r:link]", namespaces=NS):
            for attr in (f"{{{NS['r']}}}embed", f"{{{NS['r']}}}link"):
                val = node.get(attr)
                if val and val not in ids:
                    ids.append(val)
        return ids
    except TypeError:
        try:
            raw = ET.fromstring(shape.element.xml.encode("utf-8"))
            vals = []
            for elem in raw.iter():
                for attr in (f"{{{NS['r']}}}embed", f"{{{NS['r']}}}link"):
                    val = elem.attrib.get(attr)
                    if val and val not in vals:
                        vals.append(val)
            return vals
        except Exception:
            return []
    except Exception:
        return []


def extract_text_from_shape(shape: Any) -> str:
    if getattr(shape, "has_text_frame", False):
        parts: list[str] = []
        for para in shape.text_frame.paragraphs:
            line_parts = []
            for run in para.runs:
                line_parts.append(run.text)
            line = "".join(line_parts) if line_parts else para.text
            if line.strip():
                parts.append(line)
        return clean_text("\n".join(parts))
    return ""


def extract_table(shape: Any) -> list[list[str]] | None:
    if not getattr(shape, "has_table", False):
        return None
    rows: list[list[str]] = []
    for row in shape.table.rows:
        rows.append([clean_text(cell.text) for cell in row.cells])
    return rows


def extract_chart(shape: Any) -> dict[str, Any] | None:
    if not getattr(shape, "has_chart", False):
        return None
    chart = shape.chart
    info: dict[str, Any] = {
        "chart_type": str(chart.chart_type),
        "has_title": bool(chart.has_title),
        "title": None,
        "series": [],
    }
    try:
        if chart.has_title:
            info["title"] = clean_text(chart.chart_title.text_frame.text)
    except Exception:
        pass
    try:
        for series in chart.series:
            item: dict[str, Any] = {"name": str(series.name), "values": []}
            try:
                item["values"] = [v for v in series.values]
            except Exception:
                pass
            info["series"].append(item)
    except Exception:
        pass
    return info


def iter_shapes(shapes: Any, prefix: str = ""):
    for idx, shape in enumerate(shapes, start=1):
        path = f"{prefix}{idx}"
        yield path, shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes, prefix=f"{path}.")


def xml_texts(zf: zipfile.ZipFile, name: str) -> list[str]:
    try:
        root = ET.fromstring(zf.read(name))
    except KeyError:
        return []
    texts = [node.text or "" for node in root.findall(".//a:t", NS)]
    cleaned: list[str] = []
    for t in texts:
        t = clean_text(t)
        if t:
            cleaned.append(t)
    return cleaned


def rel_map(zf: zipfile.ZipFile, rel_name: str) -> dict[str, str]:
    try:
        root = ET.fromstring(zf.read(rel_name))
    except KeyError:
        return {}
    mapping: dict[str, str] = {}
    for rel in root.findall(".//rel:Relationship", NS):
        rid = rel.attrib.get("Id")
        target = rel.attrib.get("Target")
        typ = rel.attrib.get("Type", "").rsplit("/", 1)[-1]
        if rid and target:
            mapping[rid] = f"{typ}:{target}"
    return mapping


def slide_xml_name(slide_number: int) -> str:
    return f"ppt/slides/slide{slide_number}.xml"


def slide_rel_name(slide_number: int) -> str:
    return f"ppt/slides/_rels/slide{slide_number}.xml.rels"


def notes_xml_names(zf: zipfile.ZipFile) -> list[str]:
    return sorted(
        [n for n in zf.namelist() if n.startswith("ppt/notesSlides/notesSlide") and n.endswith(".xml")],
        key=lambda n: int(re.search(r"notesSlide(\d+)\.xml", n).group(1)) if re.search(r"notesSlide(\d+)\.xml", n) else 0,
    )


def notes_text_from_slide(slide: Any) -> str:
    try:
        if not slide.has_notes_slide:
            return ""
        return clean_text(slide.notes_slide.notes_text_frame.text)
    except Exception:
        return ""


def extract_presentation(path: Path) -> dict[str, Any]:
    prs = Presentation(str(path))
    result: dict[str, Any] = {
        "file": str(path),
        "slide_count": len(prs.slides),
        "slides": [],
    }
    with zipfile.ZipFile(path) as zf:
        notes_names = notes_xml_names(zf)
        for slide_idx, slide in enumerate(prs.slides, start=1):
            shape_records: list[dict[str, Any]] = []
            for shape_path, shape in iter_shapes(slide.shapes):
                text = extract_text_from_shape(shape)
                table = extract_table(shape)
                chart = extract_chart(shape)
                alt = shape_alt_text(shape)
                rels = shape_rel_ids(shape)
                if text or table or chart or alt or rels:
                    record = ShapeRecord(
                        path=shape_path,
                        shape_type=str(shape.shape_type),
                        name=getattr(shape, "name", ""),
                        left=getattr(shape, "left", None),
                        top=getattr(shape, "top", None),
                        width=getattr(shape, "width", None),
                        height=getattr(shape, "height", None),
                        text=text,
                        table=table,
                        chart=chart,
                        alt_text=alt,
                        rel_ids=rels or None,
                    )
                    shape_records.append(asdict(record))

            sx = slide_xml_name(slide_idx)
            srels = slide_rel_name(slide_idx)
            xml_fallback = xml_texts(zf, sx)
            note_xml = notes_names[slide_idx - 1] if slide_idx - 1 < len(notes_names) else None
            note_xml_texts = xml_texts(zf, note_xml) if note_xml else []
            media_refs = rel_map(zf, srels)
            result["slides"].append(
                {
                    "slide": slide_idx,
                    "shapes": shape_records,
                    "xml_texts": xml_fallback,
                    "notes_text": notes_text_from_slide(slide),
                    "notes_xml_texts": note_xml_texts,
                    "relationships": media_refs,
                }
            )
    return result


def markdown_for_deck(data: dict[str, Any]) -> str:
    lines: list[str] = []
    lines.append(f"# Raw Extract: {Path(data['file']).name}")
    lines.append("")
    lines.append(f"- Slides: {data['slide_count']}")
    lines.append("")
    for slide in data["slides"]:
        lines.append(f"## Slide {slide['slide']}")
        lines.append("")
        seen: set[str] = set()
        for rec in slide["shapes"]:
            text = rec.get("text") or ""
            if text:
                key = clean_text(text)
                if key and key not in seen:
                    lines.append(f"### Shape {rec['path']} {rec.get('name', '')}")
                    lines.append("")
                    lines.append(key)
                    lines.append("")
                    seen.add(key)
            table = rec.get("table")
            if table:
                lines.append(f"### Table {rec['path']} {rec.get('name', '')}")
                lines.append("")
                for row in table:
                    lines.append("- " + " | ".join(cell.replace("\n", " / ") for cell in row))
                lines.append("")
            chart = rec.get("chart")
            if chart:
                lines.append(f"### Chart {rec['path']} {rec.get('name', '')}")
                lines.append("")
                lines.append("```json")
                lines.append(json.dumps(chart, ensure_ascii=False, indent=2))
                lines.append("```")
                lines.append("")
            alt = rec.get("alt_text")
            if alt and alt not in seen:
                lines.append(f"### Alt/Text Metadata {rec['path']}")
                lines.append("")
                lines.append(alt)
                lines.append("")
                seen.add(alt)
            rels = rec.get("rel_ids")
            if rels:
                lines.append(f"### Relationships {rec['path']}")
                lines.append("")
                for rid in rels:
                    target = slide["relationships"].get(rid, "")
                    lines.append(f"- {rid}: {target}")
                lines.append("")

        xml_only = []
        existing_text = "\n".join(seen)
        for t in slide["xml_texts"]:
            if t and t not in existing_text and t not in xml_only:
                xml_only.append(t)
        if xml_only:
            lines.append("### XML fallback texts")
            lines.append("")
            for t in xml_only:
                lines.append(f"- {t}")
            lines.append("")

        notes = slide.get("notes_text") or ""
        note_xml = slide.get("notes_xml_texts") or []
        if notes:
            lines.append("### Speaker notes")
            lines.append("")
            lines.append(notes)
            lines.append("")
        if note_xml and not notes:
            lines.append("### Notes XML fallback texts")
            lines.append("")
            for t in note_xml:
                lines.append(f"- {t}")
            lines.append("")

        rels = slide.get("relationships") or {}
        media = {k: v for k, v in rels.items() if any(kind in v for kind in ("image:", "media:", "chart:", "diagramData:"))}
        if media:
            lines.append("### Slide media/diagram relationships")
            lines.append("")
            for rid, target in media.items():
                lines.append(f"- {rid}: {target}")
            lines.append("")
    return "\n".join(lines).rstrip() + "\n"


def main() -> int:
    if len(sys.argv) < 2:
        print("usage: python extract_pptx_full.py deck1.pptx [deck2.pptx ...]", file=sys.stderr)
        return 2
    out_dir = Path("extracted")
    out_dir.mkdir(exist_ok=True)
    index: list[dict[str, Any]] = []
    for arg in sys.argv[1:]:
        path = Path(arg)
        data = extract_presentation(path)
        safe = re.sub(r"[^A-Za-z0-9_.-]+", "_", path.stem)
        json_path = out_dir / f"{safe}.json"
        md_path = out_dir / f"{safe}_raw_extract.md"
        json_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        md_path.write_text(markdown_for_deck(data), encoding="utf-8")
        index.append({"file": str(path), "slides": data["slide_count"], "json": str(json_path), "markdown": str(md_path)})
        print(f"{path.name}: {data['slide_count']} slides -> {md_path}")
    (out_dir / "index.json").write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
