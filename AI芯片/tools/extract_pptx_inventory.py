from __future__ import annotations

import json
import re
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}


def natural_key(path: Path) -> tuple[int, str]:
    match = re.match(r"^(\d+)", path.name)
    if match:
        return int(match.group(1)), path.name.lower()
    return 999, path.name.lower()


def clean_text(text: str) -> str:
    text = text.replace("\u000b", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def iter_shape_text(shape):
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        parts = []
        for paragraph in shape.text_frame.paragraphs:
            line = "".join(run.text for run in paragraph.runs).strip()
            if line:
                parts.append(line)
        if parts:
            yield clean_text("\n".join(parts))

    if getattr(shape, "has_table", False) and shape.has_table:
        rows = []
        for row in shape.table.rows:
            rows.append([clean_text(cell.text) for cell in row.cells])
        if rows:
            yield {"table": rows}

    if getattr(shape, "shapes", None):
        for child in shape.shapes:
            yield from iter_shape_text(child)


def pptx_xml_text(zf: zipfile.ZipFile, name: str) -> list[str]:
    try:
        root = ET.fromstring(zf.read(name))
    except KeyError:
        return []
    except ET.ParseError:
        return []

    texts: list[str] = []
    for node in root.iter():
        if node.tag.endswith("}t") and node.text:
            value = clean_text(node.text)
            if value:
                texts.append(value)
    return texts


def notes_path_for_slide(zf: zipfile.ZipFile, slide_idx: int) -> str | None:
    rels_name = f"ppt/slides/_rels/slide{slide_idx}.xml.rels"
    try:
        rels_root = ET.fromstring(zf.read(rels_name))
    except KeyError:
        return None
    except ET.ParseError:
        return None

    for rel in rels_root.findall("rel:Relationship", NS):
        rel_type = rel.attrib.get("Type", "")
        target = rel.attrib.get("Target", "")
        if rel_type.endswith("/notesSlide") and target:
            if target.startswith("../"):
                return "ppt/" + target[3:]
            if target.startswith("/"):
                return target.lstrip("/")
            return "ppt/slides/" + target
    return None


def slide_title(items: list) -> str:
    for item in items:
        if isinstance(item, str):
            text = clean_text(item)
            if text:
                return text.splitlines()[0][:120]
    return "Untitled"


def dedupe_preserve_order(values: list[str]) -> list[str]:
    seen = set()
    out = []
    for value in values:
        value = clean_text(value)
        if not value or value in seen:
            continue
        seen.add(value)
        out.append(value)
    return out


def extract_presentation(path: Path) -> dict:
    prs = Presentation(str(path))
    deck = {
        "file": path.name,
        "slides": [],
        "slide_count": len(prs.slides),
    }

    with zipfile.ZipFile(path) as zf:
        for idx, slide in enumerate(prs.slides, start=1):
            shape_items = []
            for shape in slide.shapes:
                shape_items.extend(iter_shape_text(shape))

            xml_text = pptx_xml_text(zf, f"ppt/slides/slide{idx}.xml")
            notes_path = notes_path_for_slide(zf, idx)
            notes_text = pptx_xml_text(zf, notes_path) if notes_path else []

            shape_text_flat = []
            table_count = 0
            for item in shape_items:
                if isinstance(item, str):
                    shape_text_flat.append(item)
                elif isinstance(item, dict) and "table" in item:
                    table_count += 1
                    for row in item["table"]:
                        shape_text_flat.extend(cell for cell in row if cell)

            all_text = dedupe_preserve_order(shape_text_flat + xml_text + notes_text)
            deck["slides"].append(
                {
                    "number": idx,
                    "title": slide_title(shape_items) if shape_items else slide_title(all_text),
                    "shape_items": shape_items,
                    "xml_text": dedupe_preserve_order(xml_text),
                    "notes_text": dedupe_preserve_order(notes_text),
                    "all_text": all_text,
                    "table_count": table_count,
                    "image_count": sum(1 for shape in slide.shapes if shape.shape_type == 13),
                }
            )
    return deck


def write_markdown(deck: dict, output_path: Path) -> None:
    lines = [
        f"# {deck['file']}",
        "",
        f"- Slides: {deck['slide_count']}",
        "",
    ]

    for slide in deck["slides"]:
        lines.append(f"## Slide {slide['number']}: {slide['title']}")
        lines.append("")
        if slide["image_count"] or slide["table_count"]:
            lines.append(f"- Images: {slide['image_count']}; Tables: {slide['table_count']}")
            lines.append("")
        if slide["shape_items"]:
            lines.append("### Extracted Shape Text")
            for item in slide["shape_items"]:
                if isinstance(item, str):
                    for para in item.splitlines():
                        lines.append(f"- {para}")
                elif isinstance(item, dict) and "table" in item:
                    lines.append("")
                    rows = item["table"]
                    if rows:
                        width = max(len(row) for row in rows)
                        padded = [row + [""] * (width - len(row)) for row in rows]
                        lines.append("| " + " | ".join(padded[0]) + " |")
                        lines.append("| " + " | ".join(["---"] * width) + " |")
                        for row in padded[1:]:
                            lines.append("| " + " | ".join(row) + " |")
                    lines.append("")
            lines.append("")
        if slide["notes_text"]:
            lines.append("### Notes")
            for text in slide["notes_text"]:
                lines.append(f"- {text}")
            lines.append("")
        if not slide["shape_items"] and slide["all_text"]:
            lines.append("### XML Text")
            for text in slide["all_text"]:
                lines.append(f"- {text}")
            lines.append("")

    output_path.write_text("\n".join(lines), encoding="utf-8")


def main() -> int:
    root = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(".")
    out = Path(sys.argv[2]) if len(sys.argv) > 2 else root / "extracted"
    out.mkdir(parents=True, exist_ok=True)

    decks = []
    for pptx_path in sorted(root.glob("*.pptx"), key=natural_key):
        deck = extract_presentation(pptx_path)
        decks.append(deck)
        safe = re.sub(r"[^\w.-]+", "_", pptx_path.stem, flags=re.UNICODE)
        (out / f"{safe}.json").write_text(json.dumps(deck, ensure_ascii=False, indent=2), encoding="utf-8")
        write_markdown(deck, out / f"{safe}.md")
        print(f"extracted {pptx_path.name}: {deck['slide_count']} slides")

    summary = [
        {
            "file": deck["file"],
            "slide_count": deck["slide_count"],
            "titles": [slide["title"] for slide in deck["slides"]],
        }
        for deck in decks
    ]
    (out / "summary.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
